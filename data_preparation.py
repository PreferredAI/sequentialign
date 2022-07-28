from scipy.spatial.distance import euclidean, cosine
import scipy
import numpy as np

import json
from pptx import Presentation
from nltk.corpus import stopwords
import nltk

from sklearn.feature_extraction.text import TfidfVectorizer

stop_words = set(stopwords.words('english'))
tokenizer = nltk.RegexpTokenizer('[a-zA-Z]\w+')

MIN_WORD_LENGTH = 4
MIN_VOCAB_COUNT = 2

INF = 1.7976931348623157e+308

def timed_grid_generator(vocab_list, dn, deck_no, grid_type, fixed_len):
    slides_content = []
    chunks_content = []

    pptx_path = "./data/" + dn + "/" + str(deck_no) + "/slides.pptx"
    yt_path = "./data/" + dn + "/" + str(deck_no) + "/transcript.json"

    annotations = np.loadtxt("./data/" + dn + "/" + str(deck_no) + "/annotation.txt", dtype="str")
    duration = int(annotations[0])
    annotations = np.delete(annotations, [0])

    slide_rel = {}
    time_rel = {}
    temp_time_rel = {}
    temp_chunkcontent = []

    yt_idx = 0

    prs = Presentation(pptx_path)

    for slide in prs.slides:
        slide_content = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                word_tokens = tokenizer.tokenize(shape.text)
                words_lower = [w.lower() for w in word_tokens]
                words_in_vocab = [w for w in words_lower if w in vocab_list]
                if len(words_in_vocab) > 0:
                    slide_content.extend(words_in_vocab)
        if len(slide_content) > 0:
            slide_rel[len(slides_content)] = prs.slides.index(slide) + 1
            slides_content.append(slide_content)
        else:
            slide_rel[len(slides_content)] = prs.slides.index(slide) + 1
            slides_content.append(['lolzthisisplaceholder'])

    cell_duration = fixed_len
    if grid_type != "fixed":
        if grid_type == "pptx_len":
            cell_duration = duration / len(prs.slides)
        elif grid_type == "slide_len":
            cell_duration = duration / len(slide_rel)
    cell_start = 0.0
    temp_cellcount = 0
    while cell_start < (duration + cell_duration):
        temp_time_rel[temp_cellcount] = (cell_start, cell_start + cell_duration)
        cell_start += cell_duration
        temp_cellcount += 1
    # print(cell_duration, duration, len(temp_time_rel), len(slide_rel))

    with open(yt_path) as yt_json_file:
        data = json.load(yt_json_file)
        cue_groups = data["actions"][0]["updateEngagementPanelAction"]["content"]["transcriptRenderer"]["body"][
            "transcriptBodyRenderer"]["cueGroups"]
        chunk_words = []
        for cg in range(len(cue_groups)):
            if "simpleText" in cue_groups[cg]["transcriptCueGroupRenderer"]["cues"][0]["transcriptCueRenderer"][
                "cue"]:
                simple_text = \
                cue_groups[cg]["transcriptCueGroupRenderer"]["cues"][0]["transcriptCueRenderer"]["cue"][
                    "simpleText"]
            else:
                simple_text = ""
            time_ms = int(
                cue_groups[cg]["transcriptCueGroupRenderer"]["cues"][0]["transcriptCueRenderer"]["startOffsetMs"])
            time_s = time_ms / 1000
            word_tokens = tokenizer.tokenize(simple_text)
            words_lower = [w.lower() for w in word_tokens]
            words_in_vocab = [w for w in words_lower if w in vocab_list]
            if cg == (len(cue_groups) - 1):
                chunk_words.extend(words_in_vocab)
                temp_chunkcontent.append(chunk_words)
            elif time_s < temp_time_rel[yt_idx][1]:
                chunk_words.extend(words_in_vocab)
            else:
                while time_s > temp_time_rel[yt_idx][1]:
                    temp_chunkcontent.append(chunk_words)
                    yt_idx += 1
                    chunk_words = []
                chunk_words.extend(words_in_vocab)

    if len(temp_chunkcontent) < len(temp_time_rel):
        temp_chunkcontent.append([])

    for i in range(len(temp_chunkcontent)):
        if len(temp_chunkcontent[i]) > 0:
            # print(i, len(temp_chunkcontent))
            time_rel[len(chunks_content)] = temp_time_rel[i]
            chunks_content.append(temp_chunkcontent[i])
        else:
            time_rel[len(chunks_content)] = temp_time_rel[i]
            chunks_content.append(['lolzthisisplaceholder'])

    gs_timings = dict()
    for gs in range(len(annotations)):
        gsa = annotations[gs].split("-")
        if gsa[1] == '0' and gsa[2] == '0':
            continue
        gs_timings[int(gsa[0])] = (int(gsa[1]), int(gsa[2]))

    return slides_content, chunks_content, slide_rel, time_rel, gs_timings

def get_vocabulary(dn, deck_no):
    vocab_count = {}

    pptx_path = "./data/" + dn + "/" + str(deck_no) + "/slides.pptx"
    prs = Presentation(pptx_path)

    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                word_tokens = tokenizer.tokenize(shape.text)
                filtered_tokens = [w.lower() for w in word_tokens if not w in stop_words]
                words_extract = [w for w in filtered_tokens if len(w) > MIN_WORD_LENGTH]
                for word in words_extract:
                    if word in vocab_count:
                        vocab_count[word] += 1
                    else:
                        vocab_count[word] = 1

    return vocab_count

def or_time_grid_terms(startIndex, endIndex, grid_type="fixed", fixed_len=15):
    vocab_count = get_vocabulary(startIndex, endIndex)
    vocab_final = {key:val for key, val in vocab_count.items() if val > MIN_VOCAB_COUNT}
    vocab_list = list(vocab_final.keys())

    s_clocal, c_clocal, s_relocal, c_relocal, gs_timings = timed_grid_generator(vocab_list, startIndex, endIndex, grid_type, fixed_len)

    slide_strings = [' '.join(w) for w in s_clocal]
    chunk_strings = [' '.join(w) for w in c_clocal]
    num_slides = len(slide_strings)
    content_strings = []
    content_strings.extend(slide_strings)
    content_strings.extend(chunk_strings)

    vectorizer = TfidfVectorizer()
    tfidf_wm = vectorizer.fit_transform(content_strings)

    tfidf_array = tfidf_wm.toarray()

    s_vlocal = tfidf_array[0:num_slides]
    c_vlocal = tfidf_array[num_slides:len(tfidf_array)]

    return s_vlocal, c_vlocal, s_clocal, c_clocal, s_relocal, c_relocal, gs_timings